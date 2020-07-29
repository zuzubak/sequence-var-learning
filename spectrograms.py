import matplotlib.pyplot as plt
import matplotlib
from matplotlib import cm
from scipy.io import wavfile
import numpy as np
import evfuncs
from pydub import AudioSegment
import string
import random
from scipy import signal
from pptx import Presentation
from pptx.util import Inches
import os
import csv
import math

#modified from https://stackoverflow.com/questions/47147146/save-an-image-only-content-without-axes-or-anything-else-to-a-file-using-matl
def make_spectrogram(wav_file,directory):
    rate, data = wavfile.read(wav_file) #convert wav file to raw audio data
    fig,ax = plt.subplots(1) #create a new subplot within the matplotlib figure, to house the spectrogram
    fig.patch.set_visible(False) #get rid of white space around spectrograms -- from https://stackoverflow.com/questions/14908576/how-to-remove-frame-from-matplotlib-pyplot-figure-vs-matplotlib-figure-frame
    ax.axis('off') #turn off x (time) and y (frequency) axes
    pxx, freqs, bins, im = ax.specgram(x=data, Fs=rate, noverlap=511, NFFT=512,cmap=cm.jet, scale='dB')
    ax.set_ylim([0,10000])
    fig.set_size_inches(len(data)/12000,1) #scale the figure by the length (in time) of the audio sample being plotted, so that output images for longer samples are wider.
    return {'ax':ax,'fig':fig,'pxx':pxx}

def graph_spectrogram(wav_file,directory,ID='no_ID'):
    ax = make_spectrogram(wav_file,directory)['ax']
    fig = make_spectrogram(wav_file,directory)['fig']
    pxx = make_spectrogram(wav_file,directory)['pxx']
    plt.savefig(directory+'spectrograms/'+ID+'.png', dpi=300, frameon=False)
    plt.close(fig)

def graph_spectrogram2(wav_file,directory,ID='no_ID'):
    fs, x = wavfile.read(wav_file)
    f, t, Sxx = signal.spectrogram(x, fs)
    plt.pcolormesh(t, f[0:20], Sxx[0:20])
    plt.show()

def get_syls(directory,prepad=20,postpad=20,shuffle=True,n=50, notbatch='notbatch'):
    for fp in [directory+'snippets', directory+'spectrograms']:
        if os.path.exists(fp):
            for file in os.listdir(fp):
                os.remove(fp+'/'+file)
            os.rmdir(fp)
        os.makedirs(fp)
    notmat_lines_list = list(open(directory+notbatch)) #extract a list of labelled songfiles
    notmat_list = [directory+value[:value.index('\n')] for value in notmat_lines_list] #extract a list of labelled songfiles
    if shuffle==True: #randomize the order in which the files are processed, if the parameter "shuffle" is set to True
        random.shuffle(notmat_list)
    spectrograms_created = 0 #keeps a count of how many spectrograms have been generated
    f=0 #keeps track of how many songfiles have been processed
    for file in notmat_list:
        ndict = evfuncs.load_notmat(file) #loads variables stored in the notmat file pertaining to the current songfile
        wav_file = file.strip('.not.mat') #the wav file corresponding to the current notmat file should be the same, minus the ".not.mat" at the end
        file_audio = AudioSegment.from_wav(wav_file) #create an AudioSegment object of the relevant wave file, so that it can be sliced up.
        i=1 #keeps track of how many spectrograms have been generated for the current songfile
        for syl in ndict['labels']: #iterate through the labelled syllable tokens in the songfile
            if syl in string.ascii_letters and spectrograms_created < n: #check that the current syllable is a letter of the alphabet (so that "-" and "0" will be ignored)
                try:
                    onset=ndict['onsets'][i] #retrieve the onset time (in milliseconds) of the syllable within the file
                    offset=ndict['offsets'][i]
                    syl_snippet = file_audio[onset-prepad:offset+postpad] #create an AudioSegment that begins "prepad" ms before the syllable begins, and ends "postpad" milliseconds after it ends, where those two parameters are specified by the user.
                    ID = str(f)+'_'+str(i) #create an arbitrary ID for the current syllable, by concatenating the number of songs that have been processed and the number of syllables processed in the current song, with an underscore in between.
                    snippet_filename = directory+'snippets/'+ID+'.wav'
                    syl_snippet.export(snippet_filename,format='wav')
                    graph_spectrogram(snippet_filename,directory,ID=ID)
                    i+=1
                    spectrograms_created += 1
                except BaseException:
                    pass
        f+=1

        
def make_pptx(directory,fname='for_categorization.pptx', border_width_inches=0.01, buffer_width_inches=0.01, n=None):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[6] #blank slide
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = prs.slides[0].shapes
    left = top = Inches(border_width_inches)
    spectrograms_dir = directory + 'spectrograms/'
    if n == None:
        n = len(os.listdir(spectrograms_dir))
    files_list = [file for file in os.listdir(spectrograms_dir) if 'png' in file][:n]
    random.shuffle(files_list)
    i=0
    for spectrogram in files_list:
        spec_fp = spectrograms_dir + spectrogram
        if 'png' in spec_fp:
            try:
                picture = shapes.add_picture(spec_fp, left, top)
                image_width = Inches(picture.image.size[0]/picture.image.dpi[0])
                image_height = Inches(picture.image.size[1]/picture.image.dpi[1])
                #left += Inches(buffer_width_inches)
                if left + 2*image_width > prs.slide_width:
                    if top + image_height < prs.slide_height: #Create new line if adding image to current line will exceed slide borders.
                        left = Inches(border_width_inches)
                        top += image_height
                    else:
                        break
                else:
                    left += image_width
                i+=1
            except BaseException:
                pass
    prs.save(directory+fname)
    print(fname)

def batch_slides(directory):
    for bird_directory in os.listdir(directory):
        if '.' not in bird_directory:
            get_syls(directory+bird_directory+'/')
            make_pptx(directory+bird_directory+'/')

def sample_size(directory, ss_list = [30,50,75,100,150]):
    for bird_directory in os.listdir(directory):
        if '.' not in bird_directory:
            get_syls(directory+bird_directory+'/', n=max(ss_list))
            for ss in ss_list:
                make_pptx(directory+bird_directory+'/', fname=str(ss)+'_categorization.pptx', n=ss)

def sample_size_one_bird(directory, ss_list = [30,50,75,100,150]):
    get_syls(directory+'/', n=max(ss_list))
    for ss in ss_list:
        make_pptx(directory, fname=str(ss)+'_categorization.pptx', n=ss)
    
def batch_slides(directory):
    for bird_directory in os.listdir(directory):
        if '.' not in bird_directory:
            for file in os.listdir(directory+bird_directory):
                if '_categorization' in file:
                    os.rename(directory+bird_directory+'/'+file,directory+'/'+'red/'+bird_directory+'_'+file)

def get_syls2(directory,prepad=20,postpad=20,shuffle=True,n=50, notbatch='notbatch'):
    for fp in [directory+'snippets', directory+'spectrograms']:
        if os.path.exists(fp):
            for file in os.listdir(fp):
                os.remove(fp+'/'+file)
            os.rmdir(fp)
        os.makedirs(fp)
    notmat_lines_list = list(open(directory+notbatch)) #extract a list of labelled songfiles
    notmat_list = [directory+value[:value.index('\n')] for value in notmat_lines_list] #extract a list of labelled songfiles
    if shuffle==True: #randomize the order in which the files are processed, if the parameter "shuffle" is set to True
        random.shuffle(notmat_list)
    spectrograms_created = 0 #keeps a count of how many spectrograms have been generated
    f=0 #keeps track of how many songfiles have been processed
    for file in notmat_list:
        ndict = evfuncs.load_notmat(file) #loads variables stored in the notmat file pertaining to the current songfile
        wav_file = file.strip('.not.mat') #the wav file corresponding to the current notmat file should be the same, minus the ".not.mat" at the end
        file_audio = AudioSegment.from_wav(wav_file) #create an AudioSegment object of the relevant wave file, so that it can be sliced up.
        i=1 #keeps track of how many spectrograms have been generated for the current songfile
        for syl in ndict['labels']: #iterate through the labelled syllable tokens in the songfile
            if spectrograms_created < n:
                try:
                    onset=ndict['onsets'][i] #retrieve the onset time (in milliseconds) of the syllable within the file
                    offset=ndict['offsets'][i]
                    syl_snippet = file_audio[onset-prepad:offset+postpad] #create an AudioSegment that begins "prepad" ms before the syllable begins, and ends "postpad" milliseconds after it ends, where those two parameters are specified by the user.
                    ID = str(f)+'_'+str(i) #create an arbitrary ID for the current syllable, by concatenating the number of songs that have been processed and the number of syllables processed in the current song, with an underscore in between.
                    snippet_filename = directory+'snippets/'+ID+'.wav'
                    syl_snippet.export(snippet_filename,format='wav')
                    graph_spectrogram2(snippet_filename,directory,ID=ID)
                    i+=1
                    spectrograms_created += 1
                except BaseException:
                    pass
        f+=1



def latin_square(file):
    with open('./data/lat_sq.csv') as f:
        data = list(csv.reader(f))
    print(file)
    try:
        column = data[0].index(file)
        return (data[4][column],data[3][column])
    except:
        pass

def latin_square_rename(directory):
    for file in os.listdir(directory):
        if '_categorization' in file and '~' not in file:
            name = latin_square(file)
            try:
                os.rename(directory+'/'+file,directory+'rater_'+name[1]+'/'+name[0]+'.pptx')
            except:
                pass

def count_clusters(directory):
    meta_list=[]
    files_list = os.listdir(directory)
    files_list.sort
    for file in files_list:
        if '~' not in file and 'DS_Store' not in file:
            prs = Presentation(directory+file)
            out_list = [file]
            for i in range(1,len(prs.slides)):
                cluster_size = len(prs.slides[i].shapes)
                out_list.append(cluster_size)
            meta_list.append(out_list)
    with open(directory+'cluster_sizes.csv','w+') as csvfile:
        csv_writer=csv.writer(csvfile)
        for row in meta_list:
            csv_writer.writerow(row)

def entropy(csv_file):
    with open(csv_file) as f:
        data = list(csv.reader(f))
    out_list=[]
    for row in data:
        row = [int(value) for value in row if value not in ['0','']]
        probs_list = [value/sum(row) for value in row]
        print(row)
        entropy = -sum([probability * math.log(probability, 2) for probability in probs_list])
        out_list.append([str(entropy)])
    with open('./output/entropy_values.csv','w+') as csvfile:
        csv_writer=csv.writer(csvfile)
        for row in out_list:
            csv_writer.writerow(row)
    